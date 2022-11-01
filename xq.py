#!/usr/bin/python

### MODULES ###
#import datetime
#import os
#import pprint
#
#import clipboard
import codecs
import os.path
import sys
from datetime import datetime
#
from datsun import *
from qenv3 import *
import befehl
import regi
from pptx import Presentation
from pptx.shapes.group import GroupShape
#import xb
#import xt
#import xz
import kbench
#
#1:battery,2:pip,3:mygen,4:myopus

### VARIABLES ###
kbench.KBDEBUG = True
DEBUG = True
DEBUG = False
#
version = '2021-05-27'

#-------------------------------------------------
#2021-05-27
#https://github.com/shakiyam

def log(message):
	timestamp = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
	print(f'{timestamp} {message}', file=sys.stderr)

def treatShape(shape):
	lines = []
	if shape.has_text_frame:
		for paragraph in shape.text_frame.paragraphs:
			stripped = paragraph.text.strip()
			if stripped:
				lines.append(stripped)
	elif shape.has_table:
		for row in shape.table.rows:
			for cell in row.cells:
				stripped = cell.text.strip()
				if stripped:
					lines.append(stripped)
	elif isinstance(shape, GroupShape):
		for item in shape.shapes:
			lines += treatShape(item)
	return lines

#log(f'ppt2txt - version {version} by Shinichi Akiyama')

#-------------------------------------------------

##########################
### POWERPOINT zu TEXT ###
##########################
def ppt2txt(eingabe,output):
	base, ext = os.path.splitext(eingabe)
	texteingabe = f'{base}.txt'

	presentation = Presentation(eingabe)
#	log(f'{eingabe} was opened.')

	lines = []
	for i, slide in enumerate(presentation.slides):
		lines.append(f'--- Slide {i + 1} ---')
		for shape in slide.shapes:
			lines += treatShape(shape)

	with codecs.open(output, 'w', 'utf-8') as f:
		for line in lines:
			print(line, file=f)

#

##### DIREKT ###############
if __name__=='__main__':
	ppt = 'F:/gegen/nexs/2020-08-25_cvc/2_NEXS_JPN.pptx'
	ppt2txt(ppt,labomi+'x.txt')
	kbench.enfin()
